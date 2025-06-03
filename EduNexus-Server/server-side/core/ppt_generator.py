import os
import re
import base64
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import requests
from io import BytesIO
import markdown
from bs4 import BeautifulSoup
import random
import time
import ast

from api.gemini_client import GeminiProvider


class PptGenerator:
    def __init__(self):
        """
        Initialize the PPT Generator with an empty constructor.
        All parameters will be passed through the specific methods.
        """
        self.gemini_client = GeminiProvider()
        pass
            
    def _process_markdown(self, text):
        """Convert markdown to HTML and then extract text with formatting hints."""
        html = markdown.markdown(text)
        soup = BeautifulSoup(html, 'html.parser')
        
        # Extract headings
        headings = []
        for h in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
            headings.append(h.get_text())
            h.decompose()  # Remove heading from soup after extracting
            
        # Get remaining content
        content = soup.get_text().strip()
        
        # Split content into paragraphs
        paragraphs = [p for p in content.split('\n') if p.strip()]
        
        return headings[0] if headings else None, paragraphs
        
    def _add_image_to_slide(self, slide, img_src, position, text_area_bottom=None, slide_height=None):
        """
        Add an image to the slide from URL, file path or base64.
        Ensures images don't overlap with text by positioning them below text area.
        
        Args:
            slide: The slide to add the image to
            img_src: Source of the image (URL, path, or base64)
            position: Tuple of (left, top, width) for image position
            text_area_bottom: Bottom boundary of text area to prevent overlap
            slide_height: Height of the slide for calculations
        """
        try:
            left, top, width = position
            
            # Adjust top position if text_area_bottom is provided to avoid overlap
            if text_area_bottom and top < text_area_bottom:
                top = text_area_bottom + Inches(0.2)  # Add margin between text and image
            
            # Ensure image will fit on slide
            if slide_height:
                max_height = (slide_height - top) * 0.95  # 95% of remaining space
            else:
                max_height = Inches(3)  # Default max height if slide_height not provided
            
            # Handle different image source types
            if img_src.startswith('http'):
                # Web URL
                response = requests.get(img_src)
                img_stream = BytesIO(response.content)
                slide.shapes.add_picture(img_stream, left, top, width=width, height=max_height)
            elif img_src.startswith('data:image'):
                # Base64 encoded image
                img_data = re.sub('^data:image/.+;base64,', '', img_src)
                binary_data = base64.b64decode(img_data)
                img_stream = BytesIO(binary_data)
                slide.shapes.add_picture(img_stream, left, top, width=width, height=max_height)
            else:
                # Local file path
                slide.shapes.add_picture(img_src, left, top, width=width, height=max_height)
            return True
        except Exception as e:
            print(f"Error adding image: {e}")
            return False
            
    def create_title_slide(self, prs, course_name, lesson_name):
        """
        Create the title slide with course and lesson name.
        
        Args:
            prs: Presentation object
            course_name: Name of the course
            lesson_name: Name of the lesson
            
        Returns:
            The created slide
        """
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = course_name
        for paragraph in title.text_frame.paragraphs:
            paragraph.font.size = Pt(44)
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.CENTER
        
        # Enable word wrapping
        title.text_frame.word_wrap = True
        
        # Set subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = lesson_name
        for paragraph in subtitle.text_frame.paragraphs:
            paragraph.font.size = Pt(32)
            paragraph.alignment = PP_ALIGN.CENTER
        
        # Enable word wrapping for subtitle
        subtitle.text_frame.word_wrap = True
        
        return slide
        
    def create_content_slides(self, prs, markdown_list, markdown_images_list):
        """
        Process markdown content and create slides.
        
        Args:
            prs: Presentation object
            markdown_list: List of strings containing markdown text
            markdown_images_list: List of lists of strings with image paths/urls/base64
            
        Returns:
            List of created slides
        """
        slides = []
        slide_height = prs.slide_height
        
        for i, md_content in enumerate(markdown_list):
            # Get heading and content from markdown
            heading, paragraphs = self._process_markdown(md_content)
            
            if not heading:
                heading = f"Slide {i+1}"
                
            if not paragraphs:
                continue  # Skip if no content
                
            # Create slide
            content_slide_layout = prs.slide_layouts[1]  # Content with title layout
            slide = prs.slides.add_slide(content_slide_layout)
            slides.append(slide)
            
            # Add title with word wrap
            title = slide.shapes.title
            title.text = heading
            title.text_frame.word_wrap = True
            for paragraph in title.text_frame.paragraphs:
                paragraph.font.size = Pt(36)
            
            # Add content with word wrap
            text_shape = None
            content_area_bottom = Inches(2.5)  # Default bottom of text area (after title)
            
            if 1 in slide.placeholders:  # Check if there's a content placeholder
                text_shape = slide.placeholders[1]
                text_shape.text = ""  # Clear default text
                
                for para in paragraphs:
                    p = text_shape.text_frame.add_paragraph()
                    p.text = para
                    p.font.size = Pt(18)
                    p.space_after = Pt(12)  # Space between paragraphs
                
                # Enable word wrapping
                text_shape.text_frame.word_wrap = True
                
                # Calculate bottom of text area for image positioning
                if text_shape.text_frame.paragraphs:
                    # Estimate text height based on content
                    text_height = len(paragraphs) * Pt(30)  # Approximate height per paragraph
                    content_area_bottom = text_shape.top + text_height
            else:
                # If no placeholder, create a text box
                left = Inches(1)
                top = Inches(1.5)
                width = Inches(8)
                height = Inches(2.5)
                
                text_box = slide.shapes.add_textbox(left, top, width, height)
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                
                for para in paragraphs:
                    p = text_frame.add_paragraph()
                    p.text = para
                    p.font.size = Pt(18)
                    p.space_after = Pt(12)
                
                content_area_bottom = top + height
            
            # Get images for this slide (if available)
            available_images = []
            if i < len(markdown_images_list):
                available_images = markdown_images_list[i]
            
            # Add up to 2 images if available, positioning them to avoid text overlap
            if available_images:
                num_images = min(2, len(available_images))
                selected_images = random.sample(available_images, num_images)
                
                # Define positions for images (initial positions that will be adjusted)
                if num_images == 1:
                    # Center single image
                    positions = [
                        (Inches(3), content_area_bottom + Inches(0.5), Inches(4))  # Centered image
                    ]
                else:
                    # Left and right for two images
                    positions = [
                        (Inches(1), content_area_bottom + Inches(0.5), Inches(3)),  # Left image
                        (Inches(5), content_area_bottom + Inches(0.5), Inches(3))   # Right image
                    ]
                
                for j, img in enumerate(selected_images):
                    if j < len(positions):
                        self._add_image_to_slide(slide, img, positions[j], content_area_bottom, slide_height)
        
        return slides
    
    def generate_ppt(self, markdown_list, markdown_images_list, course_name, lesson_name, output_folder="generated_ppts"):
        """
        Generate the complete PowerPoint presentation.
        
        Args:
            markdown_list: List of strings containing markdown text
            markdown_images_list: List of lists of strings with image paths/urls/base64
            course_name: Name of the course
            lesson_name: Name of the lesson
            output_folder: Folder to save the presentation in (default "generated_ppts")
            
        Returns:
            Path to the saved presentation
        """
        # Create a new presentation
        prs = Presentation()
        
        # Ensure output folder exists
        if not os.path.exists(output_folder):
            os.makedirs(output_folder)
        
        # Create title slide
        self.create_title_slide(prs, course_name, lesson_name)
        
        # Create content slides
        self.create_content_slides(prs, markdown_list, markdown_images_list)
        
        # Save the presentation
        output_path = os.path.join(output_folder, f"{course_name}_{lesson_name}.pptx")
        prs.save(output_path)
        return output_path


    def generate_ppt_content(self, markdown_list : list[dict]):
        prompt = """You are a skilled and creative expert in designing PowerPoint presentations. Your task is to take the content of a course and craft an engaging and concise presentation which is content-rich.\nFor each topic in the course, you will design slides explaining the topics provided. For each topic in the course, consolidate information into fewer slides (10-15 slides) by grouping related subtopics or concepts. Aim to minimize the total number of slides while ensuring each slide contains detailed and comprehensive content that thoroughly explains the topic. Strive for a balance between clarity and depth, presenting all critical information within a logical structure. \nEnsure that the slides are:\n - Organized logically for easy understanding.\n - Engaging and aligned with professional presentation standards.\n\n # Course Content:\n\n"""
        for index, markdown in enumerate(markdown_list):
            for key, value in markdown.items():
                prompt += f"## Topic {index + 1}: {key}\n## Content {index +1}: {value}\n\n"
        
        output_format_prompt = """-------------------\n# OUTPUT FORMAT EXAMPLE: \n 
        ["# Introduction to Python\n\nPython is a high-level programming language known for its readability. It's widely used in:\n\n- Data Science\n- Web Development\n- Artificial Intelligence",
        "## Variables and Data Types\n\nPython has several built-in data types:\n\n- Strings: `\"Hello World\"`\n- Numbers: `42, 3.14`\n- Lists: `[1, 2, 3]`\n- Dictionaries: `{\"key\": \"value\"}`"
        ]
        """
    
        prompt += output_format_prompt
        while True:
            try:
                output = self.gemini_client.generate_response(prompt, remove_literals=False)
                pattern = r"\[\s*(.*)\s*\]"
                match = re.search(pattern, output, re.DOTALL)
                extracted_list_string = match.group(0)
                extracted_list = ast.literal_eval(extracted_list_string)
                return extracted_list
            except Exception as e:
                print(f"Exception while creating PPT:{e}.\nTrying again in 10 seconds...")
                time.sleep(10)
            