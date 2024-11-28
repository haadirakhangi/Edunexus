from api.gemini_client import GeminiProvider
import markdown
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import base64
import io
import requests
import re
import ast
import time

# class PptGenerator:
#     def __init__(self):
#         self.gemini_client = GeminiProvider()
#         self.prs = Presentation()
#         self.blank_slide_layout = self.prs.slide_layouts[6]

#     def convert_markdown_to_formatted_content(self, markdown_text):
#         html = markdown.markdown(markdown_text, extensions=['tables'])
#         soup = BeautifulSoup(html, 'html.parser')

#         parsed_content = {'paragraphs': [], 'tables': []}

#         for element in soup.find_all(['h1', 'h2', 'h3', 'p', 'ul', 'table']):
#             if element.name in ['h1', 'h2', 'h3']:
#                 header_level = int(element.name[1])
#                 parsed_content['paragraphs'].append({
#                     'text': element.get_text(),
#                     'style': 'header',
#                     'level': header_level
#                 })
#             elif element.name == 'p':
#                 parsed_content['paragraphs'].append({
#                     'text': element.get_text(),
#                     'style': 'normal',
#                     'level': 0
#                 })
#             elif element.name == 'ul':
#                 for li in element.find_all('li'):
#                     parsed_content['paragraphs'].append({
#                         'text': li.get_text(),
#                         'style': 'bullet',
#                         'level': 1
#                     })
#             elif element.name == 'table':
#                 table_data = {'headers': [], 'rows': []}
#                 headers = element.find_all('th')
#                 if headers:
#                     table_data['headers'] = [header.get_text() for header in headers]

#                 rows = element.find_all('tr')[1:] if headers else element.find_all('tr')
#                 for row in rows:
#                     row_data = [cell.get_text() for cell in row.find_all(['td', 'th'])]
#                     table_data['rows'].append(row_data)

#                 parsed_content['tables'].append(table_data)
#         return parsed_content

#     def add_title(self, slide, title):
#         left, top, width, height = Inches(0.5), Inches(0.5), Inches(9), Inches(1)
#         txBox = slide.shapes.add_textbox(left, top, width, height)
#         tf = txBox.text_frame
#         p = tf.paragraphs[0]
#         p.text = title
#         p.alignment = PP_ALIGN.CENTER
#         font = p.font
#         font.size = Pt(32)
#         font.color.rgb = RGBColor(0, 51, 102)
#         font.bold = True

#     def add_content(self, slide, markdown_content):
#         content = self.convert_markdown_to_formatted_content(markdown_content)
#         left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(4)
#         txBox = slide.shapes.add_textbox(left, top, width, height)
#         tf = txBox.text_frame

#         for item in content['paragraphs']:
#             p = tf.add_paragraph()
#             p.text = item['text']
#             if item['style'] == 'header':
#                 p.font.size = Pt(24 - (item['level'] * 4))
#                 p.font.bold = True
#                 p.level = item['level'] - 1
#             elif item['style'] == 'bullet':
#                 p.level = 1
#                 p.font.size = Pt(16)
#             else:
#                 p.font.size = Pt(16)

#         for table_data in content['tables']:
#             self.add_table_to_slide(slide, table_data)

#     def add_table_to_slide(self, slide, table_data):
#         left, top, width, height = Inches(0.5), Inches(5), Inches(9), Inches(2)
#         num_columns = len(table_data['headers']) if table_data['headers'] else len(table_data['rows'][0])
#         num_rows = len(table_data['rows']) + (1 if table_data['headers'] else 0)
#         table = slide.shapes.add_table(num_rows, num_columns, left, top, width, height).table

#         if table_data['headers']:
#             for col, header in enumerate(table_data['headers']):
#                 cell = table.cell(0, col)
#                 cell.text = header
#                 cell.text_frame.paragraphs[0].font.bold = True

#         start_row = 1 if table_data['headers'] else 0
#         for row_idx, row_data in enumerate(table_data['rows'], start=start_row):
#             for col_idx, cell_data in enumerate(row_data):
#                 table.cell(row_idx, col_idx).text = str(cell_data)

#     def add_image_to_slide(self, slide, image_source, left, top, width, height):
#         def is_base64(data):
#             try:
#                 if not data.isprintable() or any(c.isspace() for c in data):
#                     return False
#                 decoded = base64.b64decode(data, validate=True)
#                 return base64.b64encode(decoded).decode() == data
#             except Exception:
#                 return False

#         try:
#             if is_base64(image_source):
#                 image_bytes = base64.b64decode(image_source.split(",")[1] if "," in image_source else image_source)
#                 slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, width, height)
#             elif image_source.startswith(("http:", "https:")):
#                 response = requests.get(image_source, stream=True)
#                 if response.status_code == 200:
#                     slide.shapes.add_picture(io.BytesIO(response.content), left, top, width, height)
#             else:
#                 with Image.open(image_source):  # Validation
#                     slide.shapes.add_picture(image_source, left, top, width, height)
#         except Exception as e:
#             print(f"Error adding image: {e}")

#     def add_images(self, slide, images):
#         max_width, max_height = Inches(10), Inches(5.5)
#         if len(images) == 1:
#             img_width, img_height = Inches(6), Inches(4)
#             left, top = (max_width - img_width) / 2, Inches(3)
#             self.add_image_to_slide(slide, images[0], left, top, img_width, img_height)
#         elif len(images) == 2:
#             img_width, img_height = Inches(4), Inches(3)
#             positions = [(Inches(1), Inches(3)), (Inches(5.5), Inches(3))]
#             for img, pos in zip(images, positions):
#                 self.add_image_to_slide(slide, img, pos[0], pos[1], img_width, img_height)
#         elif len(images) >= 3:
#             img_width, img_height = Inches(2.5), Inches(2)
#             positions = [
#                 (Inches(1), Inches(2)), (Inches(6.5), Inches(2)),
#                 (Inches(1), Inches(4.5)), (Inches(6.5), Inches(4.5))
#             ]
#             for img, pos in zip(images[:4], positions):
#                 self.add_image_to_slide(slide, img, pos[0], pos[1], img_width, img_height)

#     def create_presentation(self, slide_data_list : list[dict], output_filename):
#         for slide_data in slide_data_list:
#             slide = self.prs.slides.add_slide(self.blank_slide_layout)
#             self.add_title(slide, slide_data.get('title', 'Untitled'))
#             self.add_content(slide, slide_data.get('slide_content', ''))
#             self.add_images(slide, slide_data.get('images', []))

#         self.prs.save(output_filename)
#         print(f"Presentation saved as {output_filename}")

#     def generate_ppt_content(self, markdown_list : list[dict]):
#         prompt = """You are a skilled and creative expert in designing PowerPoint presentations. Your task is to take the content of a course and summarize it effectively to create an engaging and concise presentation.\nFor each topic in the course, you will design slides summarizing the material provided. The number of slides (or dictionaries in the list) can be as much as necessary to adequately explain the topics. However, the number of slides must not be less than the number of topics and can be greater if required for clarity and depth. If a topic contains subtopics or sections, ensure that each subtopic is covered on separate slides where appropriate.\nEnsure that the summaries are:\n - Comprehensive yet concise.\n - Organized logically for easy understanding.\n - Engaging and aligned with professional presentation standards.\n - Tailored to the target audience of the presentation, in this case students, to match their level of expertise and interest.\n\n # Course Content:\n\n"""
#         for index, markdown in enumerate(markdown_list):
#             for key, value in markdown.items():
#                 prompt += f"## Topic {index + 1}: {key}\n## Content {index +1}: {value}\n\n"

#         output_format_prompt = """-------------------\n# OUTPUT FORMAT INSTRUCTIONS:\nThe output should be a list of dictionaries where each dictionary represents the content of one slide. Each dictionary within the list must include two keys:\n'title': The title of the slide, summarizing the main point of the topic.\n'slide_content': The concise content to be added to the slide in the form of a markdowns, expressed in bullet points (5-7) or short paragraphs for clarity.\nProvide the output strictly as a list of dictionaries following the format as described. Do not include any additional text or explanation outside of the given structure. Do not include unnecessary text like ```json in the output. It should strictly only be a list of dictionary."""
#         prompt += output_format_prompt
#         output = self.gemini_client.generate_response(prompt, remove_literals=False)
#         pattern = r"\[\s*(.*)\s*\]"
#         match = re.search(pattern, output, re.DOTALL)
#         extracted_list_string = match.group(0)
#         extracted_list = ast.literal_eval(extracted_list_string)
#         return extracted_list


class PptGenerator:
    def __init__(self):
        self.gemini_client = GeminiProvider()

    @staticmethod
    def parse_html_content(html_text):
        soup = BeautifulSoup(html_text, 'html.parser')

        parsed_content = []

        for element in soup.find_all(['h1', 'h2', 'h3', 'p', 'ul', 'li', 'b', 'i', 'u']):
            if element.name in ['h1', 'h2', 'h3']:
                header_level = int(element.name[1])
                parsed_content.append({
                    'text': element.get_text(),
                    'style': 'header',
                    'level': header_level
                })
            elif element.name == 'p':
                parsed_content.append({
                    'text': element.get_text(),
                    'style': 'normal',
                    'level': 0
                })
            elif element.name == 'ul':
                for li in element.find_all('li'):
                    parsed_content.append({
                        'text': li.get_text(),
                        'style': 'bullet',
                        'level': 1
                    })
            elif element.name in ['b', 'i', 'u']:
                # For formatting inside paragraphs
                parent = element.find_parent(['p', 'li'])
                if parent:
                    parsed_content.append({
                        'text': parent.get_text(),
                        'style': 'formatted',
                        'format': element.name
                    })

        return parsed_content

    @staticmethod
    def add_title(slide, title):
        left, top, width, height = Inches(0.5), Inches(0.5), Inches(9), Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        font = p.font
        font.size = Pt(32)
        font.color.rgb = RGBColor(0, 51, 102)
        font.bold = True

    @staticmethod
    def add_content(slide, html_content):
        content = PptGenerator.parse_html_content(html_content)
        left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        for item in content:
            if 'text' in item:
                # Add paragraph for each text block and apply formatting
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT  # Align text to the left
                p.text = ""  # Initialize text to avoid duplicate issues
                PptGenerator.add_text_with_formatting(p, item['text'])

                # Adjust font size and bullet style based on type
                if item['style'] == 'header':
                    p.font.size = Pt(24 - (item['level'] * 4))
                    p.font.bold = True
                    p.level = item['level'] - 1
                elif item['style'] == 'bullet':
                    p.level = 1
                    p.font.size = Pt(18)
                    p.text = "• " + item['text']  # Ensure bullet symbol is added manually
                elif item['style'] == 'normal':
                    p.font.size = Pt(18)

                # Check if text exceeds slide size and reduce font size if necessary
                while PptGenerator.check_text_exceeds_bounds(p):
                    p.font.size -= Pt(1)
                    if p.font.size < Pt(10):
                        break
                    # Check if text exceeds slide size and reduce font size if necessary
        while PptGenerator.check_text_exceeds_bounds(p):
            p.font.size -= Pt(1)
            if p.font.size < Pt(10):
                break

    @staticmethod
    def add_text_with_formatting(paragraph, html_text):
        """ Parse HTML and add text with inline formatting (bold, italic, underline). """
        soup = BeautifulSoup(html_text, 'html.parser')

        for element in soup.descendants:
            if isinstance(element, str):
                # Add text directly if it's a string
                run = paragraph.add_run()
                run.text = element
            elif element.name in ['b', 'i', 'u']:
                run = paragraph.add_run()
                run.text = element.get_text()

                # Apply formatting based on the tag
                if element.name == 'b':
                    run.font.bold = True
                if element.name == 'i':
                    run.font.italic = True
                if element.name == 'u':
                    run.font.underline = True



    @staticmethod
    def check_text_exceeds_bounds(paragraph):
        max_characters_per_line = 80
        max_lines = 10

        text_length = len(paragraph.text)
        lines = (text_length // max_characters_per_line) + 1
        return lines > max_lines

    @staticmethod
    def create_presentation(slide_data_list : list[dict], output_filename):
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        for slide_data in slide_data_list:
            slide = prs.slides.add_slide(blank_slide_layout)
            PptGenerator.add_title(slide, slide_data.get('title', 'Untitled'))
            PptGenerator.add_content(slide, slide_data.get('slide_content', ''))

        prs.save(output_filename)
        print(f"Presentation saved as {output_filename}")
    
    def generate_ppt_content(self, markdown_list : list[dict]):
        prompt = """You are a skilled and creative expert in designing PowerPoint presentations. Your task is to take the content of a course and craft an engaging and concise presentation which is content-rich.\nFor each topic in the course, you will design slides explaining the topics provided. For each topic in the course, consolidate information into fewer slides (10-15 slides) by grouping related subtopics or concepts. Aim to minimize the total number of slides while ensuring each slide contains detailed and comprehensive content that thoroughly explains the topic. Strive for a balance between clarity and depth, presenting all critical information within a logical structure. \nEnsure that the slides are:\n - Organized logically for easy understanding.\n - Engaging and aligned with professional presentation standards.\n\n # Course Content:\n\n"""
        for index, markdown in enumerate(markdown_list):
            for key, value in markdown.items():
                prompt += f"## Topic {index + 1}: {key}\n## Content {index +1}: {value}\n\n"

        output_format_prompt = """-------------------\n# OUTPUT FORMAT INSTRUCTIONS:\nThe output should be a list of dictionaries where each dictionary represents the content of one slide. Each dictionary within the list must include two keys:\n'title': The title of the slide, summarizing the main point of the topic.\n'slide_content': Detailed and comprehensive content to include on the slide, formatted as HTML. This should contain:\n - Paragraphs with detailed explanations, covering all relevant information.\n - Bullet points that expand on key details, ensuring clarity and depth.\n - Do not include any heading tags in the slide_content.\nProvide the output strictly as a list of dictionaries following the format as described. Do not include unnecessary text like ```json in the output. It should strictly only be a list of dictionary."""
        prompt += output_format_prompt
        while True:
            try:
                output = self.gemini_client.generate_response(prompt, remove_literals=False)
                pattern = r"\[\s*(.*)\s*\]"
                match = re.search(pattern, output, re.DOTALL)
                extracted_list_string = match.group(0)
                extracted_list = ast.literal_eval(extracted_list_string)
                return extracted_list
            except Exception as e:
                print(f"Exception while creating PPT:{e}.\nTrying again in 10 seconds...")
                time.sleep(10)
