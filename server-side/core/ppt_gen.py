import os
import re
import markdown
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image
import base64
import io
import requests
from pptx.util import Inches
from PIL import Image


def convert_markdown_to_formatted_content(markdown_text):
    """
    Convert markdown to a structured format for PowerPoint

    Returns:
    {
        'paragraphs': [
            {'text': 'string', 'style': 'header/normal/bold', 'level': 0-3},
        ],
        'tables': [
            {'headers': [], 'rows': []}
        ]
    }
    """
    html = markdown.markdown(markdown_text, extensions=['tables'])
    soup = BeautifulSoup(html, 'html.parser')

    parsed_content = {
        'paragraphs': [],
        'tables': []
    }

    for element in soup.find_all(['h1', 'h2', 'h3', 'p', 'ul', 'table']):
        if element.name in ['h1', 'h2', 'h3']:
            # Headers
            header_level = int(element.name[1])
            parsed_content['paragraphs'].append({
                'text': element.get_text(),
                'style': 'header',
                'level': header_level
            })

        elif element.name == 'p':
            parsed_content['paragraphs'].append({
                'text': element.get_text(),
                'style': 'normal',
                'level': 0
            })

        elif element.name == 'ul':
            for li in element.find_all('li'):
                parsed_content['paragraphs'].append({
                    'text': li.get_text(),
                    'style': 'bullet',
                    'level': 1
                })

        elif element.name == 'table':
            table_data = {
                'headers': [],
                'rows': []
            }

            headers = element.find_all('th')
            if headers:
                table_data['headers'] = [header.get_text() for header in headers]

            rows = element.find_all('tr')[1:] if headers else element.find_all('tr')
            for row in rows:
                row_data = [cell.get_text() for cell in row.find_all(['td', 'th'])]
                table_data['rows'].append(row_data)

            parsed_content['tables'].append(table_data)

    return parsed_content

def create_powerpoint_from_slides(slide_data_list, output_filename='presentation.pptx'):
    """
    Create a PowerPoint presentation from a list of slide dictionaries.

    Each slide dictionary should have:
    - title: String for slide title
    - content: Markdown-formatted content string
    - images: List of image file paths
    """
    prs = Presentation()

    blank_slide_layout = prs.slide_layouts[6]

    def add_title(slide, title):
        """Add a title to the slide with custom formatting"""
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER

        # Formatting
        font = p.font
        font.size = Pt(32)
        font.color.rgb = RGBColor(0, 51, 102)  
        font.bold = True

    def add_content(slide, markdown_content):
        """Add formatted content to the slide"""
        content = convert_markdown_to_formatted_content(markdown_content)

        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(4)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame


        for item in content['paragraphs']:
            p = tf.add_paragraph()
            p.text = item['text']

            if item['style'] == 'header':
                p.font.size = Pt(24 - (item['level'] * 4))
                p.font.bold = True
                p.level = item['level'] - 1
            elif item['style'] == 'bullet':
                p.level = 1  
                p.font.size = Pt(16)
            else:
                p.font.size = Pt(16)

        if content['tables']:
            for table_data in content['tables']:
                add_table_to_slide(slide, table_data)

    def add_table_to_slide(slide, table_data):
        """
        Add a table to the slide

        Args:
            slide: PowerPoint slide object
            table_data: Dictionary with headers and rows
        """
        left = Inches(0.5)
        top = Inches(5)
        width = Inches(9)
        height = Inches(2)

        num_columns = len(table_data['headers']) if table_data['headers'] else len(table_data['rows'][0])
        num_rows = len(table_data['rows']) + (1 if table_data['headers'] else 0)

        shapes = slide.shapes
        table = shapes.add_table(
            num_rows,
            num_columns,
            left,
            top,
            width,
            height
        ).table

        if table_data['headers']:
            for col, header in enumerate(table_data['headers']):
                cell = table.cell(0, col)
                cell.text = header
                cell.text_frame.paragraphs[0].font.bold = True

        start_row = 1 if table_data['headers'] else 0
        for row_idx, row_data in enumerate(table_data['rows'], start=start_row):
            for col_idx, cell_data in enumerate(row_data):
                table.cell(row_idx, col_idx).text = str(cell_data)


    def add_image_to_slide(slide, image_source, left, top, width, height):
        """Add a single image to the slide with support for Base64, URLs, and local paths."""

        def is_base64(data):
          """Check if the string is Base64-encoded."""
          try:
              if not data.isprintable() or any(c.isspace() for c in data):
                  return False

              decoded = base64.b64decode(data, validate=True)
              return base64.b64encode(decoded).decode() == data
          except Exception:
              return False

        if is_base64(image_source):  
            try:
                if "," in image_source:
                    image_source = image_source.split(",")[1] 
                image_bytes = base64.b64decode(image_source)
                image_stream = io.BytesIO(image_bytes)
                slide.shapes.add_picture(image_stream, left, top, width, height)
            except Exception as e:
                print(f"Error adding Base64 image: {e}")

        elif image_source.startswith("http:") or image_source.startswith("https:"):  
            try:
                response = requests.get(image_source, stream=True)
                print("response")
                if response.status_code == 200:
                    image_stream = io.BytesIO(response.content)
                    slide.shapes.add_picture(image_stream, left, top, width, height)
                else:
                    print(f"Error fetching image from URL: HTTP {response.status_code}")
            except Exception as e:
                print(f"Error adding URL image: {e}")

        else: 
            try:
                with Image.open(image_source):  
                    slide.shapes.add_picture(image_source, left, top, width, height)
            except FileNotFoundError:
                print(f"Image not found: {image_source}")
            except Exception as e:
                print(f"Error adding image {image_source}: {e}")


    def add_images(slide, images):
        """Add images to the slide with dynamic sizing and positioning."""
        max_width = Inches(10)
        max_height = Inches(5.5)

        if len(images) == 1:
            img_width = Inches(6)
            img_height = Inches(4)
            left = (max_width - img_width) / 2
            top = Inches(3)
            add_image_to_slide(slide, images[0], left, top, img_width, img_height)

        elif len(images) == 2:
            img_width = Inches(4)
            img_height = Inches(3)
            left1 = Inches(1)
            left2 = Inches(5.5)
            top = Inches(3)
            add_image_to_slide(slide, images[0], left1, top, img_width, img_height)
            add_image_to_slide(slide, images[1], left2, top, img_width, img_height)

        elif len(images) == 3:
            img_width = Inches(3)
            img_height = Inches(2)
            top_img_left = (max_width - img_width) / 2
            bottom_left1 = Inches(1.5)
            bottom_left2 = Inches(6)

            add_image_to_slide(slide, images[0], top_img_left, Inches(2), img_width, img_height)
            add_image_to_slide(slide, images[1], bottom_left1, Inches(4), img_width, img_height)
            add_image_to_slide(slide, images[2], bottom_left2, Inches(4), img_width, img_height)

        elif len(images) >= 4:
            img_width = Inches(2.5)
            img_height = Inches(2)
            positions = [
                (Inches(1), Inches(2)),    
                (Inches(6.5), Inches(2)),  
                (Inches(1), Inches(4.5)),  
                (Inches(6.5), Inches(4.5)) 
            ]

            for i in range(min(4, len(images))):
                add_image_to_slide(slide, images[i], positions[i][0], positions[i][1], img_width, img_height)



    for slide_data in slide_data_list:
        slide = prs.slides.add_slide(blank_slide_layout)
        add_title(slide, slide_data.get('title', 'Untitled'))
        add_content(slide, slide_data.get('content', ''))
        add_images(slide, slide_data.get('images', []))

    prs.save(output_filename)
    print(f"Presentation saved as {output_filename}")

